"""
PowerPoint presentation generation module.
Creates .pptx files with screenshots and documentation.
"""

import os
import logging
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image


class PresentationGenerator:
    """Generates PowerPoint presentations from user flow data."""
    
    def __init__(self, config: Dict[str, Any]):
        """Initialize the presentation generator."""
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.prs = None
        
    def create_presentation(self) -> None:
        """Create a new PowerPoint presentation."""
        self.prs = Presentation()
        
        # Set slide dimensions
        slide_width = self.config.get('PRESENTATION', {}).get('slide_width', 9144000)
        slide_height = self.config.get('PRESENTATION', {}).get('slide_height', 6858000)
        
        self.prs.slide_width = slide_width
        self.prs.slide_height = slide_height
        
        self.logger.info("Created new PowerPoint presentation")
    
    def add_cover_slide(self, title: str, subtitle: str = "") -> None:
        """Add a cover slide to the presentation."""
        if not self.prs:
            self.create_presentation()
        
        # Use title slide layout
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Set subtitle if provided
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            
            # Format subtitle
            subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
            subtitle_paragraph.font.size = Pt(24)
            subtitle_paragraph.font.color.rgb = RGBColor(102, 102, 102)  # Gray
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
        self.logger.info("Added cover slide")
    
    def add_table_of_contents_slide(self, steps: List[Dict[str, Any]]) -> None:
        """Add a table of contents slide."""
        if not self.prs:
            self.create_presentation()
        
        # Use blank layout
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Table of Contents"
        
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add steps list
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        content_frame = content_box.text_frame
        content_frame.clear()
        
        for i, step in enumerate(steps, 1):
            paragraph = content_frame.add_paragraph()
            paragraph.text = f"{i}. {step.get('name', 'Unknown Step')}"
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
            paragraph.space_after = Pt(12)
        
        self.logger.info("Added table of contents slide")
    
    def add_step_slide(self, step: Dict[str, Any], screenshot_path: str, step_number: int) -> None:
        """Add a slide for a single step with screenshot."""
        if not self.prs:
            self.create_presentation()
        
        # Use blank layout
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add step title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = f"Step {step_number}: {step.get('name', 'Unknown Step')}"
        
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(24)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        
        # Add screenshot if it exists
        if os.path.exists(screenshot_path):
            try:
                # Get image dimensions
                with Image.open(screenshot_path) as img:
                    img_width, img_height = img.size
                
                # Calculate scaling to fit in slide
                max_width = Inches(8)
                max_height = Inches(5)
                
                scale_x = max_width / Inches(img_width / 96)  # Assuming 96 DPI
                scale_y = max_height / Inches(img_height / 96)
                scale = min(scale_x, scale_y, 1.0)  # Don't scale up
                
                scaled_width = Inches(img_width / 96) * scale
                scaled_height = Inches(img_height / 96) * scale
                
                # Center the image
                left = Inches(5) - scaled_width / 2
                top = Inches(2.5) - scaled_height / 2
                
                slide.shapes.add_picture(screenshot_path, left, top, scaled_width, scaled_height)
                self.logger.info(f"Added screenshot to slide: {screenshot_path}")
                
            except Exception as e:
                self.logger.error(f"Failed to add screenshot {screenshot_path}: {e}")
                # Add placeholder text instead
                placeholder_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                placeholder_frame = placeholder_box.text_frame
                placeholder_frame.text = f"Screenshot not available: {screenshot_path}"
                placeholder_paragraph = placeholder_frame.paragraphs[0]
                placeholder_paragraph.font.size = Pt(16)
                placeholder_paragraph.font.color.rgb = RGBColor(150, 150, 150)
                placeholder_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add step description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
        desc_frame = desc_box.text_frame
        desc_frame.text = step.get('description', 'No description provided')
        
        desc_paragraph = desc_frame.paragraphs[0]
        desc_paragraph.font.size = Pt(16)
        desc_paragraph.font.color.rgb = RGBColor(51, 51, 51)
        desc_paragraph.alignment = PP_ALIGN.LEFT
        
        self.logger.info(f"Added step slide: {step.get('name', 'Unknown Step')}")
    
    def add_summary_slide(self, steps: List[Dict[str, Any]]) -> None:
        """Add a summary slide."""
        if not self.prs:
            self.create_presentation()
        
        # Use blank layout
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Summary"
        
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add summary content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        content_frame = content_box.text_frame
        content_frame.clear()
        
        # Add flow overview
        overview_paragraph = content_frame.add_paragraph()
        overview_paragraph.text = f"This user flow consists of {len(steps)} steps:"
        overview_paragraph.font.size = Pt(18)
        overview_paragraph.font.bold = True
        overview_paragraph.font.color.rgb = RGBColor(51, 51, 51)
        overview_paragraph.space_after = Pt(12)
        
        # Add step summary
        for i, step in enumerate(steps, 1):
            step_paragraph = content_frame.add_paragraph()
            step_paragraph.text = f"• Step {i}: {step.get('name', 'Unknown Step')}"
            step_paragraph.font.size = Pt(16)
            step_paragraph.font.color.rgb = RGBColor(51, 51, 51)
            step_paragraph.space_after = Pt(8)
        
        # Add completion note
        completion_paragraph = content_frame.add_paragraph()
        completion_paragraph.text = "\nThis documentation was automatically generated by CaptureBot."
        completion_paragraph.font.size = Pt(14)
        completion_paragraph.font.color.rgb = RGBColor(102, 102, 102)
        completion_paragraph.font.italic = True
        completion_paragraph.alignment = PP_ALIGN.CENTER
        
        self.logger.info("Added summary slide")
    
    def save_presentation(self, filename: str = "user_flow_documentation.pptx") -> str:
        """Save the presentation to a file."""
        if not self.prs:
            raise ValueError("No presentation to save")
        
        # Ensure the filename has .pptx extension
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        self.prs.save(filename)
        self.logger.info(f"Presentation saved: {filename}")
        return filename
    
    def generate_presentation(self, steps: List[Dict[str, Any]], screenshots: List[str], 
                            title: str = "User Flow Documentation") -> str:
        """Generate a complete presentation from steps and screenshots."""
        self.create_presentation()
        
        # Add cover slide
        self.add_cover_slide(title, f"Automated documentation of {len(steps)} steps")
        
        # Add table of contents
        self.add_table_of_contents_slide(steps)
        
        # Add step slides
        for i, (step, screenshot) in enumerate(zip(steps, screenshots), 1):
            self.add_step_slide(step, screenshot, i)
        
        # Add summary slide
        self.add_summary_slide(steps)
        
        # Save presentation
        filename = self.save_presentation()
        return filename
